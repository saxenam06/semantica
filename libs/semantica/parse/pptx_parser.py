"""
PowerPoint document parser for Semantica framework.

This module handles PPTX file parsing using python-pptx
for presentation content extraction.
"""

from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from pptx import Presentation

from ..utils.exceptions import ProcessingError, ValidationError
from ..utils.logging import get_logger


@dataclass
class SlideContent:
    """Slide content representation."""
    
    slide_number: int
    title: str
    text: str
    notes: str = ""
    shapes: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)


@dataclass
class PPTXData:
    """PPTX presentation representation."""
    
    title: str
    slides: List[SlideContent]
    metadata: Dict[str, Any] = field(default_factory=dict)


class PPTXParser:
    """PPTX document parser."""
    
    def __init__(self, **config):
        """
        Initialize PPTX parser.
        
        Args:
            **config: Parser configuration
        """
        self.logger = get_logger("pptx_parser")
        self.config = config
    
    def parse(self, file_path: Union[str, Path], **options) -> PPTXData:
        """
        Parse PPTX file.
        
        Args:
            file_path: Path to PPTX file
            **options: Parsing options:
                - extract_images: Whether to extract images (default: False)
                - extract_notes: Whether to extract speaker notes (default: True)
                - extract_shapes: Whether to extract shape information (default: False)
                
        Returns:
            PPTXData: Parsed presentation data
        """
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise ValidationError(f"PPTX file not found: {file_path}")
        
        if not file_path.suffix.lower() == '.pptx':
            raise ValidationError(f"File is not a PPTX: {file_path}")
        
        try:
            prs = Presentation(str(file_path))
            
            # Extract metadata
            metadata = self._extract_metadata(prs)
            
            # Extract slides
            slides = []
            for idx, slide in enumerate(prs.slides, 1):
                slide_content = self._parse_slide(slide, idx, options)
                slides.append(slide_content)
            
            # Get presentation title (from first slide or core properties)
            title = metadata.get("title", f"Presentation_{file_path.stem}")
            if slides and slides[0].title:
                title = slides[0].title
            
            return PPTXData(
                title=title,
                slides=slides,
                metadata={
                    **metadata,
                    "total_slides": len(slides),
                    "file_path": str(file_path)
                }
            )
            
        except Exception as e:
            self.logger.error(f"Failed to parse PPTX {file_path}: {e}")
            raise ProcessingError(f"Failed to parse PPTX: {e}")
    
    def extract_text(self, file_path: Union[str, Path]) -> str:
        """
        Extract text from PPTX.
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            str: Extracted text from all slides
        """
        result = self.parse(file_path, extract_images=False, extract_shapes=False)
        return "\n\n".join([f"Slide {s.slide_number}: {s.title}\n{s.text}" for s in result.slides])
    
    def extract_slides(self, file_path: Union[str, Path], **options) -> List[SlideContent]:
        """
        Extract slides from PPTX.
        
        Args:
            file_path: Path to PPTX file
            **options: Parsing options
            
        Returns:
            list: List of slide contents
        """
        result = self.parse(file_path, **options)
        return result.slides
    
    def _parse_slide(self, slide, slide_number: int, options: Dict[str, Any]) -> SlideContent:
        """Parse individual slide."""
        title = ""
        text_parts = []
        notes = ""
        shapes = []
        images = []
        
        # Extract title and text from shapes
        for shape in slide.shapes:
            # Extract title
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                if shape.placeholder_format.idx == 0:  # Title placeholder
                    if hasattr(shape, "text"):
                        title = shape.text.strip()
            
            # Extract text from text boxes
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            
            # Extract images
            if options.get("extract_images", False) and hasattr(shape, "image"):
                try:
                    image_info = {
                        "width": shape.width,
                        "height": shape.height,
                        "format": shape.image.ext
                    }
                    images.append(image_info)
                except Exception:
                    pass
            
            # Extract shape information
            if options.get("extract_shapes", False):
                shape_info = {
                    "type": shape.shape_type,
                    "has_text": hasattr(shape, "text")
                }
                shapes.append(shape_info)
        
        # Extract notes
        if options.get("extract_notes", True) and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if hasattr(notes_slide, "notes_text_frame"):
                notes = notes_slide.notes_text_frame.text
        
        # Get title from first text part if not found
        if not title and text_parts:
            title = text_parts[0]
            text_parts = text_parts[1:]
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            text="\n".join(text_parts),
            notes=notes,
            shapes=shapes,
            images=images
        )
    
    def _extract_metadata(self, prs: Presentation) -> Dict[str, Any]:
        """Extract presentation metadata."""
        metadata = {}
        
        try:
            core_props = prs.core_properties
            metadata = {
                "title": core_props.title,
                "author": core_props.author,
                "subject": core_props.subject,
                "keywords": core_props.keywords,
                "comments": core_props.comments,
                "created": str(core_props.created) if core_props.created else None,
                "modified": str(core_props.modified) if core_props.modified else None,
                "last_modified_by": core_props.last_modified_by
            }
        except Exception:
            pass
        
        return metadata
